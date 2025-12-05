from pptx import Presentation
import os

# Create a dummy pptx
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx is working!"

prs.save('test.pptx')

# Read it back using file_extractor logic
try:
    prs = Presentation('test.pptx')
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    extracted = "\n".join(text)
    print(f"Extracted: {extracted}")
    
    if "Hello, World!" in extracted:
        print("SUCCESS")
    else:
        print("FAILURE")
        
except Exception as e:
    print(f"Error: {e}")
    
finally:
    if os.path.exists('test.pptx'):
        os.remove('test.pptx')
