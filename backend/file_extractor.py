import os
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

def extract_text(file_path: str) -> str:
    """
    Extracts text content from various file formats.
    """
    if not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".docx":
            if Document is None:
                return "Error: python-docx library not installed."
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif ext == ".pptx":
            if Presentation is None:
                return "Error: python-pptx library not installed."
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        elif ext == ".html" or ext == ".htm":
            if BeautifulSoup is None:
                return "Error: beautifulsoup4 library not installed."
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "lxml")
                return soup.get_text(separator="\n")
        
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        
        else:
            return f"Unsupported file format: {ext}"
            
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""
